"""rag.chunking のファイル読み込み（PDF / Office）と分割ヘルパのテスト。優先度: RAG。"""

from pathlib import Path

import pytest
from pypdf import PdfWriter
from rag.chunking import (
    _split_sections,
    _split_sysml_code_blocks,
    read_excel_file,
    read_pdf_file,
    read_powerpoint_file,
    read_word_file,
)


def test_read_pdf_file_minimal_pdf(tmp_path: Path):
    pdf_path = tmp_path / "one.pdf"
    w = PdfWriter()
    w.add_blank_page(width=72, height=72)
    with open(pdf_path, "wb") as f:
        w.write(f)

    text, pages = read_pdf_file(pdf_path)
    assert len(pages) >= 1
    assert pages[0][0] == 1
    assert isinstance(text, str)


def test_read_word_file_roundtrip(tmp_path: Path):
    pytest.importorskip("docx")
    from docx import Document

    doc = Document()
    doc.add_paragraph("RAG chunking test")
    p = tmp_path / "w.docx"
    doc.save(p)

    full, paras = read_word_file(p)
    assert "RAG" in full
    assert len(paras) >= 1


def test_read_word_file_missing_raises(tmp_path: Path):
    pytest.importorskip("docx")

    with pytest.raises(FileNotFoundError):
        read_word_file(tmp_path / "missing.docx")


def test_read_excel_file_sheet(tmp_path: Path):
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws["A1"] = "alpha"
    ws["B1"] = "beta"
    p = tmp_path / "book.xlsx"
    wb.save(p)

    full, sheets = read_excel_file(p)
    assert "alpha" in full
    assert len(sheets) >= 1


def test_read_powerpoint_title_slide(tmp_path: Path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TitleLine"
    p = tmp_path / "deck.pptx"
    prs.save(p)

    full, slides = read_powerpoint_file(p)
    assert "TitleLine" in full
    assert slides and "TitleLine" in slides[0][1]


def test_split_sysml_code_blocks_empty_lines():
    assert _split_sysml_code_blocks("") == [("text", "")]


def test_split_sections_single_block_without_heading():
    """見出し無しは 1 セクション。"""
    secs = _split_sections("only body\nsecond line")
    assert len(secs) == 1
    assert secs[0][2]  # non-empty body


def test_split_sections_numeric_headings():
    secs = _split_sections("1.1 TitleA\n\nbody here\n\n2 Next\n\nmore")
    assert len(secs) >= 2
