"""ファイル読み込みとチャンク分割モジュール。

このモジュールは、テキストファイル、Markdownファイル、PDFファイル、
Word、Excel、PowerPointファイルの読み込みと、テキストのチャンク分割機能を提供します。
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Iterable, List, Optional, Tuple

from pypdf import PdfReader

# Word/Excel/PowerPoint用のライブラリ（オプションインポート）
try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from openpyxl import load_workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


CHUNK_SIZE = 1000  # 文字数ベース
CHUNK_OVERLAP = 100

_NUM_HEADING_RE = re.compile(r"^\s*(\d+(?:\.\d+)*)\s+(.+?)\s*$")
_MD_HEADING_RE = re.compile(r"^\s*(#{1,6})\s+(.+?)\s*$")

_SYSML_START_RE = re.compile(
    r"^\s*(package\b|part\s+def\b|action\s+def\b|activity\s+def\b|type\s+def\b|"
    r"connection\s+def\b|requirement\s+def\b|view\b|import\b|private\s+import\b)\s*",
    re.IGNORECASE,
)


def _split_sections(text: str) -> List[Tuple[Optional[str], Optional[str], str]]:
    """章/節（見出し）でテキストを分割する（最小実装）。

    - 数字見出し: "8.1 Title" のような行
    - Markdown見出し: "# Title", "## Title" のような行
    - 見出しが無ければ全体を1セクションとして返す
    """
    lines = (text or "").splitlines()
    if not lines:
        return [(None, None, "")]

    # Markdown見出しの場合はレベルに応じて連番IDを作る（例: 1, 1.1, 1.1.1）
    counters = [0, 0, 0, 0, 0, 0]
    current_id: Optional[str] = None
    current_title: Optional[str] = None
    buf: List[str] = []
    out: List[Tuple[Optional[str], Optional[str], str]] = []

    def flush() -> None:
        nonlocal buf, out, current_id, current_title
        if buf:
            out.append((current_id, current_title, "\n".join(buf).strip()))
            buf = []

    for line in lines:
        m_num = _NUM_HEADING_RE.match(line)
        m_md = _MD_HEADING_RE.match(line)

        if m_num:
            flush()
            current_id = m_num.group(1)
            current_title = m_num.group(2)
            continue

        if m_md:
            flush()
            level = len(m_md.group(1))
            title = m_md.group(2)
            # 更新: levelカウンタ++、下位をリセット
            counters[level - 1] += 1
            for i in range(level, len(counters)):
                counters[i] = 0
            parts = [str(c) for c in counters[:level] if c > 0]
            current_id = ".".join(parts) if parts else None
            current_title = title
            continue

        buf.append(line)

    flush()
    if not out:
        return [(None, None, (text or "").strip())]
    return out


def _split_sysml_code_blocks(text: str) -> List[Tuple[str, str]]:
    """SysMLコードっぽいブロックを「塊」として抽出する（最小ヒューリスティック）。

    Returns:
        List[Tuple[kind, segment_text]] where kind is "text" or "sysml_code".
    """
    lines = (text or "").splitlines()
    if not lines:
        return [("text", "")]

    segments: List[Tuple[str, str]] = []
    buf: List[str] = []
    code: List[str] = []
    in_code = False
    brace_depth = 0

    def flush_text() -> None:
        nonlocal buf
        if buf:
            segments.append(("text", "\n".join(buf).strip()))
            buf = []

    def flush_code() -> None:
        nonlocal code, brace_depth
        if code:
            segments.append(("sysml_code", "\n".join(code).strip()))
            code = []
        brace_depth = 0

    for i, line in enumerate(lines):
        is_start = _SYSML_START_RE.match(line) is not None
        if not in_code and is_start:
            flush_text()
            in_code = True

        if in_code:
            code.append(line)
            # ブレース深さで雑にブロック終端を推定
            brace_depth += line.count("{")
            brace_depth -= line.count("}")
            if brace_depth <= 0:
                # 末尾がセミコロン/閉じブレースで、次行が空か見出しなら「塊」として閉じる
                next_line = lines[i + 1] if i + 1 < len(lines) else ""
                next_is_heading = bool(_NUM_HEADING_RE.match(next_line) or _MD_HEADING_RE.match(next_line))
                if line.strip().endswith((";", "}")) and (not next_line.strip() or next_is_heading):
                    in_code = False
                    flush_code()
            continue

        buf.append(line)

    if in_code:
        flush_code()
    else:
        flush_text()

    # 空セグメント除去（ただし順序は保つ）
    cleaned: List[Tuple[str, str]] = [(k, t) for (k, t) in segments if t.strip()]
    return cleaned if cleaned else [("text", (text or "").strip())]


def chunk_text_semantic(text: str) -> List[Tuple[str, dict]]:
    """意味単位（章/節）を壊さずにチャンク化し、メタデータも返す（ステップ2）。

    - 章/節を検出してセクション単位に分割
    - セクション内で SysMLコードっぽいブロックは 1チャンクにまとめる
    - それ以外は従来の sentence+1000文字チャンクを使用

    Returns:
        List[Tuple[chunk_text, meta]]
        meta: {"section_id","section_title","chunk_kind","code_language"}
    """
    out: List[Tuple[str, dict]] = []
    for section_id, section_title, section_text in _split_sections(text):
        if not section_text:
            continue
        for kind, seg in _split_sysml_code_blocks(section_text):
            if not seg.strip():
                continue
            if kind == "sysml_code":
                out.append(
                    (
                        seg,
                        {
                            "section_id": section_id,
                            "section_title": section_title,
                            "chunk_kind": "sysml_code",
                            "code_language": "sysml",
                        },
                    )
                )
            else:
                for c in chunk_text(seg):
                    out.append(
                        (
                            c,
                            {
                                "section_id": section_id,
                                "section_title": section_title,
                                "chunk_kind": "text",
                            },
                        )
                    )
    return out


def read_text_file(path: Path) -> str:
    """テキストファイルをUTF-8エンコーディングで読み込む。

    Args:
        path: 読み込むファイルのパス。

    Returns:
        str: ファイルの内容。

    Raises:
        FileNotFoundError: ファイルが存在しない場合。
        UnicodeDecodeError: ファイルがUTF-8でデコードできない場合。
    """
    return path.read_text(encoding="utf-8")


def read_markdown_file(path: Path) -> str:
    """MarkdownファイルをUTF-8エンコーディングで読み込む。

    現時点ではテキストファイルと同様に扱う。

    Args:
        path: 読み込むMarkdownファイルのパス。

    Returns:
        str: ファイルの内容。

    Raises:
        FileNotFoundError: ファイルが存在しない場合。
        UnicodeDecodeError: ファイルがUTF-8でデコードできない場合。
    """
    return path.read_text(encoding="utf-8")


def read_pdf_file(path: Path) -> Tuple[str, List[Tuple[int, str]]]:
    """PDFファイルからテキストを抽出する。

    Args:
        path: 読み込むPDFファイルのパス。

    Returns:
        Tuple[str, List[Tuple[int, str]]]: 
            - 全ページを連結したテキスト
            - 各ページの(ページ番号, ページテキスト)のリスト

    Raises:
        FileNotFoundError: ファイルが存在しない場合。
        Exception: PDFの読み込みに失敗した場合。

    Example:
        >>> text, pages = read_pdf_file(Path("document.pdf"))
        >>> print(f"Total pages: {len(pages)}")
        Total pages: 10
    """
    reader = PdfReader(str(path))
    pages: List[Tuple[int, str]] = []
    texts: List[str] = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        pages.append((i + 1, text))
        texts.append(text)
    return "\n".join(texts), pages


def read_word_file(path: Path) -> Tuple[str, List[Tuple[int, str]]]:
    """Wordファイル（.docx）からテキストを抽出する。

    Args:
        path: 読み込むWordファイルのパス。

    Returns:
        Tuple[str, List[Tuple[int, str]]]:
            - 全段落を連結したテキスト
            - 各段落の(段落番号, 段落テキスト)のリスト

    Raises:
        FileNotFoundError: ファイルが存在しない場合。
        ImportError: python-docxがインストールされていない場合。
        Exception: Wordファイルの読み込みに失敗した場合。

    Example:
        >>> text, paragraphs = read_word_file(Path("document.docx"))
        >>> print(f"Total paragraphs: {len(paragraphs)}")
        Total paragraphs: 10
    """
    if not HAS_DOCX:
        raise ImportError(
            "python-docx がインストールされていません。"
            "Wordファイルを読み込むには `pip install python-docx` を実行してください。"
        )
    if not path.exists():
        raise FileNotFoundError(f"ファイルが存在しません: {path}")
    doc = DocxDocument(str(path))
    paragraphs: List[Tuple[int, str]] = []
    texts: List[str] = []
    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if text:  # 空の段落はスキップ
            paragraphs.append((i + 1, text))
            texts.append(text)
    return "\n".join(texts), paragraphs


def read_excel_file(path: Path) -> Tuple[str, List[Tuple[int, str]]]:
    """Excelファイル（.xlsx, .xls）からテキストを抽出する。

    Args:
        path: 読み込むExcelファイルのパス。

    Returns:
        Tuple[str, List[Tuple[int, str]]]:
            - 全シート・全セルを連結したテキスト
            - 各シートの(シート番号, シートテキスト)のリスト

    Raises:
        FileNotFoundError: ファイルが存在しない場合。
        ImportError: openpyxlがインストールされていない場合。
        Exception: Excelファイルの読み込みに失敗した場合。

    Example:
        >>> text, sheets = read_excel_file(Path("document.xlsx"))
        >>> print(f"Total sheets: {len(sheets)}")
        Total sheets: 3
    """
    if not HAS_OPENPYXL:
        raise ImportError(
            "openpyxl がインストールされていません。"
            "Excelファイルを読み込むには `pip install openpyxl` を実行してください。"
        )
    workbook = load_workbook(str(path), data_only=True)
    sheets: List[Tuple[int, str]] = []
    all_texts: List[str] = []
    
    for sheet_idx, sheet in enumerate(workbook.worksheets):
        sheet_texts: List[str] = []
        for row in sheet.iter_rows(values_only=True):
            row_texts: List[str] = []
            for cell_value in row:
                if cell_value is not None:
                    # 数値や日付も文字列に変換
                    cell_str = str(cell_value).strip()
                    if cell_str:
                        row_texts.append(cell_str)
            if row_texts:
                sheet_texts.append(" | ".join(row_texts))
        
        if sheet_texts:
            sheet_text = "\n".join(sheet_texts)
            sheets.append((sheet_idx + 1, sheet_text))
            all_texts.append(sheet_text)
    
    workbook.close()
    return "\n\n".join(all_texts), sheets


def read_powerpoint_file(path: Path) -> Tuple[str, List[Tuple[int, str]]]:
    """PowerPointファイル（.pptx）からテキストを抽出する。

    Args:
        path: 読み込むPowerPointファイルのパス。

    Returns:
        Tuple[str, List[Tuple[int, str]]]:
            - 全スライドを連結したテキスト
            - 各スライドの(スライド番号, スライドテキスト)のリスト

    Raises:
        FileNotFoundError: ファイルが存在しない場合。
        ImportError: python-pptxがインストールされていない場合。
        Exception: PowerPointファイルの読み込みに失敗した場合。

    Example:
        >>> text, slides = read_powerpoint_file(Path("presentation.pptx"))
        >>> print(f"Total slides: {len(slides)}")
        Total slides: 15
    """
    if not HAS_PPTX:
        raise ImportError(
            "python-pptx がインストールされていません。"
            "PowerPointファイルを読み込むには `pip install python-pptx` を実行してください。"
        )
    if not path.exists():
        raise FileNotFoundError(f"ファイルが存在しません: {path}")
    prs = Presentation(str(path))
    slides: List[Tuple[int, str]] = []
    texts: List[str] = []
    
    for i, slide in enumerate(prs.slides):
        slide_texts: List[str] = []
        # スライド内のすべてのテキストフレームからテキストを抽出
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    slide_texts.append(text)
        
        if slide_texts:
            slide_text = "\n".join(slide_texts)
            slides.append((i + 1, slide_text))
            texts.append(slide_text)
    
    return "\n\n".join(texts), slides


def split_to_sentences(text: str) -> List[str]:
    """テキストをセンテンス単位に分割する。

    日本語/英語混在を想定した簡易的なセンテンス分割を実行します。
    句点（。、.）、感嘆符（！、!）、疑問符（？、?）、改行で分割します。

    Args:
        text: 分割対象のテキスト。

    Returns:
        List[str]: センテンスのリスト。空のセンテンスは除外されます。

    Example:
        >>> text = "これは最初の文です。これは2番目の文です。"
        >>> sentences = split_to_sentences(text)
        >>> print(sentences)
        ['これは最初の文です。', 'これは2番目の文です。']
    """
    # 日本語の句点・改行・英語のピリオドなどで分割
    raw_sentences = re.split(r"(?<=[。．！？\.\?\!])\s*|\n+", text)
    sentences = [s.strip() for s in raw_sentences if s.strip()]
    return sentences


def chunk_sentences(sentences: Iterable[str]) -> List[str]:
    """センテンスをまとめてチャンクを作成する。

    センテンス単位でまとめて、約 ``CHUNK_SIZE`` 文字のチャンクを作成します。
    オーバーラップは前チャンクの末尾 ``CHUNK_OVERLAP`` 文字を
    次のチャンクの先頭に付加する形で実現します。

    Args:
        sentences: センテンスのイテラブル。

    Returns:
        List[str]: チャンクのリスト。各チャンクは約 ``CHUNK_SIZE`` 文字。

    Example:
        >>> sentences = ["文1", "文2", "文3"]
        >>> chunks = chunk_sentences(sentences)
        >>> print(len(chunks))
        1
    """
    chunks: List[str] = []
    current: List[str] = []
    current_len = 0

    for sent in sentences:
        if current_len + len(sent) > CHUNK_SIZE and current:
            chunk_text = "".join(current)
            chunks.append(chunk_text)

            # オーバーラップ領域を作成
            overlap_text = chunk_text[-CHUNK_OVERLAP:] if CHUNK_OVERLAP > 0 else ""
            current = [overlap_text, sent]
            current_len = len(overlap_text) + len(sent)
        else:
            current.append(sent)
            current_len += len(sent)

    if current:
        chunks.append("".join(current))

    return chunks


def chunk_text(text: str) -> List[str]:
    """生テキストをセンテンス分割し、チャンクを生成する高レベル関数。

    Args:
        text: チャンク分割対象のテキスト。

    Returns:
        List[str]: チャンクのリスト。

    Example:
        >>> text = "長いテキストです。複数の文があります。"
        >>> chunks = chunk_text(text)
        >>> print(f"Number of chunks: {len(chunks)}")
        Number of chunks: 1
    """
    sentences = split_to_sentences(text)
    return chunk_sentences(sentences)



